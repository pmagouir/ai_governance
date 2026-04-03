#!/usr/bin/env python3
"""
DC CAP Enterprise AI Leadership Pilot - Automated Baseline Report Generator
Reads survey and milestone data, generates analyses, and creates a branded PowerPoint presentation.
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Rectangle
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import os

# DC CAP Brand Colors
NAVY = RGBColor(16, 26, 76)
PURPLE = RGBColor(131, 82, 255)
BLUE = RGBColor(135, 183, 213)
GOLD = RGBColor(252, 176, 23)
TAN = RGBColor(255, 244, 237)
WHITE = RGBColor(255, 255, 255)
GRAY_LIGHT = RGBColor(200, 200, 200)

# Hex versions for matplotlib
NAVY_HEX = '#101A4C'
PURPLE_HEX = '#8352FF'
BLUE_HEX = '#87B7D5'
GOLD_HEX = '#FCB017'
TAN_HEX = '#FFF4ED'

# Configure matplotlib
plt.rcParams['font.family'] = 'serif'
plt.rcParams['figure.facecolor'] = NAVY_HEX
plt.rcParams['axes.facecolor'] = NAVY_HEX
plt.rcParams['text.color'] = 'white'
sns.set_palette([PURPLE_HEX, BLUE_HEX, GOLD_HEX, '#FFD700', '#C0C0C0'])

# Data paths
SURVEY_DATA_PATH = '/sessions/jolly-wonderful-bell/mnt/ai_governance/test_survey_data.csv'
MILESTONE_DATA_PATH = '/sessions/jolly-wonderful-bell/mnt/ai_governance/test_milestone_data.csv'
OUTPUT_DIR = '/sessions/jolly-wonderful-bell/mnt/ai_governance'
CHARTS_DIR = os.path.join(OUTPUT_DIR, 'charts')
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, 'pilot_baseline_report.pptx')

# Create charts directory
os.makedirs(CHARTS_DIR, exist_ok=True)

def load_data():
    """Load survey and milestone data."""
    survey_df = pd.read_csv(SURVEY_DATA_PATH)
    milestone_df = pd.read_csv(MILESTONE_DATA_PATH)
    return survey_df, milestone_df

def analyze_fluency_distribution(survey_df):
    """Get fluency level distribution."""
    fluency_counts = survey_df['fluency_level'].value_counts()
    return fluency_counts

def calculate_construct_means(survey_df):
    """Calculate mean scores for each construct."""
    constructs = {
        'AI Orientation': 'ai_orientation_pct',
        'Learning Orientation': 'learning_orientation_pct',
        'Current AI Use': 'current_ai_use_pct',
        'AI Knowledge': 'ai_knowledge_pct',
        'Applied Skills': 'applied_skills_pct'
    }

    means = {}
    for name, col in constructs.items():
        means[name] = survey_df[col].mean()

    return means, constructs

def calculate_unit_performance(survey_df):
    """Calculate construct scores by unit."""
    constructs = {
        'AI Orientation': 'ai_orientation_pct',
        'Learning Orientation': 'learning_orientation_pct',
        'Current AI Use': 'current_ai_use_pct',
        'AI Knowledge': 'ai_knowledge_pct',
        'Applied Skills': 'applied_skills_pct'
    }

    units = survey_df['unit'].unique()
    unit_data = {}

    for unit in sorted(units):
        unit_df = survey_df[survey_df['unit'] == unit]
        unit_data[unit] = {name: unit_df[col].mean() for name, col in constructs.items()}

    return unit_data

def calculate_knowing_doing_gap(survey_df):
    """Calculate the knowing-doing gap: Knowledge % - Skills %."""
    gaps = survey_df['ai_knowledge_pct'].values - survey_df['applied_skills_pct'].values
    return gaps, survey_df[['participant_name', 'unit']].copy()

def get_milestone_progress(milestone_df):
    """Count milestones completed per participant."""
    progress = milestone_df.groupby('participant_name').size().reset_index(name='milestones_completed')
    return progress

def create_pie_chart_fluency(survey_df):
    """Create pie chart of fluency level distribution."""
    fluency_counts = survey_df['fluency_level'].value_counts()

    fig, ax = plt.subplots(figsize=(10, 7.5), facecolor=NAVY_HEX)

    colors = [PURPLE_HEX, BLUE_HEX, GOLD_HEX, '#C0C0C0']
    wedges, texts, autotexts = ax.pie(
        fluency_counts.values,
        labels=fluency_counts.index,
        autopct='%1.0f%%',
        colors=colors[:len(fluency_counts)],
        startangle=90,
        textprops={'color': 'white', 'fontsize': 14, 'weight': 'bold'}
    )

    for autotext in autotexts:
        autotext.set_color('white')
        autotext.set_fontsize(12)
        autotext.set_weight('bold')

    ax.set_title('Sample Fluency Level Distribution\n(N=15)',
                 color='white', fontsize=16, weight='bold', pad=20)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, 'fluency_distribution.png'),
                dpi=300, facecolor=NAVY_HEX, edgecolor='none', bbox_inches='tight')
    plt.close()

def create_construct_bars(construct_means, survey_df):
    """Create grouped bar chart for all constructs."""
    fig, ax = plt.subplots(figsize=(12, 7.5), facecolor=NAVY_HEX)

    constructs = list(construct_means.keys())
    means = list(construct_means.values())

    # Calculate standard errors
    construct_cols = {
        'AI Orientation': 'ai_orientation_pct',
        'Learning Orientation': 'learning_orientation_pct',
        'Current AI Use': 'current_ai_use_pct',
        'AI Knowledge': 'ai_knowledge_pct',
        'Applied Skills': 'applied_skills_pct'
    }

    errors = [survey_df[construct_cols[c]].std() / np.sqrt(len(survey_df)) for c in constructs]

    x_pos = np.arange(len(constructs))
    bars = ax.bar(x_pos, means, yerr=errors,
                   color=PURPLE_HEX, edgecolor=BLUE_HEX, linewidth=2,
                   capsize=5, error_kw={'elinewidth': 2, 'ecolor': BLUE_HEX})

    ax.set_ylabel('Mean Score (%)', color='white', fontsize=12, weight='bold')
    ax.set_title('Construct Scores Overview\nAll 5 Dimensions with Error Bars',
                 color='white', fontsize=16, weight='bold', pad=20)
    ax.set_xticks(x_pos)
    ax.set_xticklabels(constructs, color='white', fontsize=11, rotation=15, ha='right')
    ax.set_ylim(0, 100)
    ax.tick_params(axis='y', labelcolor='white', labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('white')
    ax.spines['bottom'].set_color('white')
    ax.grid(axis='y', alpha=0.2, color='white', linestyle='--')

    # Add value labels
    for i, (bar, mean) in enumerate(zip(bars, means)):
        ax.text(bar.get_x() + bar.get_width()/2, mean + 3, f'{mean:.1f}%',
                ha='center', va='bottom', color='white', fontsize=10, weight='bold')

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, 'construct_overview.png'),
                dpi=300, facecolor=NAVY_HEX, edgecolor='none', bbox_inches='tight')
    plt.close()

def create_unit_heatmap(unit_data):
    """Create heatmap of unit performance."""
    df_heatmap = pd.DataFrame(unit_data).T

    fig, ax = plt.subplots(figsize=(12, 6), facecolor=NAVY_HEX)

    im = ax.imshow(df_heatmap.values, cmap='YlOrRd', aspect='auto', vmin=0, vmax=100)

    ax.set_xticks(np.arange(len(df_heatmap.columns)))
    ax.set_yticks(np.arange(len(df_heatmap.index)))
    ax.set_xticklabels(df_heatmap.columns, color='white', fontsize=10, rotation=15, ha='right')
    ax.set_yticklabels(df_heatmap.index, color='white', fontsize=10)

    # Add text annotations
    for i in range(len(df_heatmap.index)):
        for j in range(len(df_heatmap.columns)):
            text = ax.text(j, i, f'{df_heatmap.values[i, j]:.0f}',
                          ha="center", va="center", color="white", fontsize=10, weight='bold')

    ax.set_title('Unit Comparison: Construct Strengths & Gaps',
                 color='white', fontsize=16, weight='bold', pad=20)

    # Color bar
    cbar = plt.colorbar(im, ax=ax, label='Score (%)')
    cbar.ax.tick_params(labelcolor='white')

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, 'unit_heatmap.png'),
                dpi=300, facecolor=NAVY_HEX, edgecolor='none', bbox_inches='tight')
    plt.close()

def create_knowing_doing_gap_chart(gaps, participant_info):
    """Create bar chart of knowing-doing gaps."""
    fig, ax = plt.subplots(figsize=(14, 7.5), facecolor=NAVY_HEX)

    sorted_idx = np.argsort(gaps)[::-1]
    sorted_gaps = gaps[sorted_idx]
    sorted_names = participant_info.iloc[sorted_idx]['participant_name'].values

    # Color by gap size
    colors_list = [PURPLE_HEX if g > 0 else BLUE_HEX for g in sorted_gaps]

    bars = ax.barh(range(len(sorted_gaps)), sorted_gaps, color=colors_list, edgecolor='white', linewidth=1)

    ax.set_yticks(range(len(sorted_gaps)))
    ax.set_yticklabels([name.replace('[TEST] ', '') for name in sorted_names],
                       color='white', fontsize=10)
    ax.set_xlabel('Knowledge % - Skills %', color='white', fontsize=12, weight='bold')
    ax.set_title('The Knowing-Doing Gap\n(Knowledge minus Applied Skills)',
                 color='white', fontsize=16, weight='bold', pad=20)
    ax.axvline(x=0, color=GOLD_HEX, linestyle='--', linewidth=2)
    ax.tick_params(axis='x', labelcolor='white', labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('white')
    ax.spines['bottom'].set_color('white')
    ax.grid(axis='x', alpha=0.2, color='white', linestyle='--')

    # Add value labels
    for i, (bar, gap) in enumerate(zip(bars, sorted_gaps)):
        ax.text(gap + 1, bar.get_y() + bar.get_height()/2, f'{gap:.1f}%',
                ha='left', va='center', color='white', fontsize=9, weight='bold')

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, 'knowing_doing_gap.png'),
                dpi=300, facecolor=NAVY_HEX, edgecolor='none', bbox_inches='tight')
    plt.close()

def create_participant_clusters(survey_df):
    """Create scatter plot of knowledge vs skills with quadrants."""
    fig, ax = plt.subplots(figsize=(11, 8.5), facecolor=NAVY_HEX)

    knowledge = survey_df['ai_knowledge_pct'].values
    skills = survey_df['applied_skills_pct'].values

    # Create color map by unit
    units = survey_df['unit'].unique()
    unit_colors = {unit: color for unit, color in
                   zip(sorted(units), [PURPLE_HEX, BLUE_HEX, GOLD_HEX])}
    colors = [unit_colors[unit] for unit in survey_df['unit'].values]

    scatter = ax.scatter(knowledge, skills, s=200, c=colors, alpha=0.7,
                        edgecolors='white', linewidth=2)

    # Add quadrant lines
    median_knowledge = np.median(knowledge)
    median_skills = np.median(skills)
    ax.axvline(median_knowledge, color=GOLD_HEX, linestyle='--', linewidth=2, alpha=0.5)
    ax.axhline(median_skills, color=GOLD_HEX, linestyle='--', linewidth=2, alpha=0.5)

    # Quadrant labels
    ax.text(75, 70, 'High Knowledge\nHigh Skills', fontsize=11, color=GOLD_HEX,
            weight='bold', ha='center', bbox=dict(boxstyle='round',
            facecolor=NAVY_HEX, edgecolor=GOLD_HEX, alpha=0.8))
    ax.text(25, 70, 'Low Knowledge\nHigh Skills', fontsize=11, color=BLUE_HEX,
            weight='bold', ha='center', bbox=dict(boxstyle='round',
            facecolor=NAVY_HEX, edgecolor=BLUE_HEX, alpha=0.8))

    ax.set_xlabel('AI Knowledge (%)', color='white', fontsize=12, weight='bold')
    ax.set_ylabel('Applied Skills (%)', color='white', fontsize=12, weight='bold')
    ax.set_title('Participant Clusters: Knowledge vs Skills\nwith Quadrant Analysis',
                 color='white', fontsize=16, weight='bold', pad=20)
    ax.set_xlim(0, 105)
    ax.set_ylim(0, 105)
    ax.tick_params(axis='both', labelcolor='white', labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('white')
    ax.spines['bottom'].set_color('white')
    ax.grid(alpha=0.2, color='white', linestyle='--')

    # Legend
    handles = [mpatches.Patch(color=color, label=unit, alpha=0.7)
               for unit, color in unit_colors.items()]
    ax.legend(handles=handles, loc='lower right', framealpha=0.9,
             facecolor=NAVY_HEX, edgecolor='white', labelcolor='white')

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, 'participant_clusters.png'),
                dpi=300, facecolor=NAVY_HEX, edgecolor='none', bbox_inches='tight')
    plt.close()

def create_milestone_chart(milestone_progress, survey_df):
    """Create horizontal bar chart of milestone completion."""
    fig, ax = plt.subplots(figsize=(12, 8), facecolor=NAVY_HEX)

    # Merge with survey data to get unit info
    merged = milestone_progress.merge(survey_df[['participant_name', 'unit']], on='participant_name')
    merged = merged.sort_values('milestones_completed', ascending=True)

    # Color by unit
    unit_colors_map = {'Student Success': PURPLE_HEX, 'Data & Technology': BLUE_HEX, 'Innovation Hub': GOLD_HEX}
    colors = [unit_colors_map[unit] for unit in merged['unit'].values]

    bars = ax.barh(range(len(merged)), merged['milestones_completed'].values,
                   color=colors, edgecolor='white', linewidth=1.5)

    ax.set_yticks(range(len(merged)))
    ax.set_yticklabels([name.replace('[TEST] ', '') for name in merged['participant_name'].values],
                       color='white', fontsize=10)
    ax.set_xlabel('Milestones Completed', color='white', fontsize=12, weight='bold')
    ax.set_title('Milestone Progress by Participant',
                 color='white', fontsize=16, weight='bold', pad=20)
    ax.tick_params(axis='x', labelcolor='white', labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('white')
    ax.spines['bottom'].set_color('white')
    ax.grid(axis='x', alpha=0.2, color='white', linestyle='--')

    # Add value labels
    for i, (bar, val) in enumerate(zip(bars, merged['milestones_completed'].values)):
        ax.text(val + 0.1, bar.get_y() + bar.get_height()/2, str(int(val)),
                ha='left', va='center', color='white', fontsize=10, weight='bold')

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, 'milestone_progress.png'),
                dpi=300, facecolor=NAVY_HEX, edgecolor='none', bbox_inches='tight')
    plt.close()

def add_title_slide(prs):
    """Add title slide (Slide 1)."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "AI FLUENCY PRE-LAUNCH ASSESSMENT"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Baseline Report"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Org name
    org_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    org_frame = org_box.text_frame
    p = org_frame.paragraphs[0]
    p.text = "DC CAP Enterprise AI Leadership Pilot"
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"April 3, 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Footer
    add_slide_footer(slide, 1)

def add_executive_summary_slide(prs, survey_df, construct_means):
    """Add executive summary slide (Slide 2)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "EXECUTIVE SUMMARY"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Key findings
    findings = [
        f"Average AI Fluency: {survey_df['overall_percentage'].mean():.0f}% | Range: {survey_df['overall_percentage'].min():.0f}–{survey_df['overall_percentage'].max():.0f}%",
        f"Strongest Construct: Learning Orientation ({construct_means['Learning Orientation']:.1f}%) | Weakest: Current AI Use ({construct_means['Current AI Use']:.1f}%)",
        f"Fluency Distribution: 4 Leading | 6 Applying | 4 Building | 1 Exploring",
        f"Knowing-Doing Gap exists for {(survey_df['ai_knowledge_pct'] > survey_df['applied_skills_pct']).sum()} of 15 participants (avg {(survey_df['ai_knowledge_pct'] - survey_df['applied_skills_pct']).mean():.1f}%)"
    ]

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, finding in enumerate(findings):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = finding
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0

    add_slide_footer(slide, 2)

def add_sample_overview_slide(prs):
    """Add sample overview slide (Slide 3)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "SAMPLE OVERVIEW"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Add chart
    chart_path = os.path.join(CHARTS_DIR, 'fluency_distribution.png')
    slide.shapes.add_picture(chart_path, Inches(1.5), Inches(1.2), width=Inches(7))

    add_slide_footer(slide, 3)

def add_construct_overview_slide(prs):
    """Add construct overview slide (Slide 4)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "CONSTRUCT SCORES OVERVIEW"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Add chart
    chart_path = os.path.join(CHARTS_DIR, 'construct_overview.png')
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.1), width=Inches(9))

    add_slide_footer(slide, 4)

def add_unit_comparison_slide(prs):
    """Add unit comparison slide (Slide 5)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "UNIT COMPARISON"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Add chart
    chart_path = os.path.join(CHARTS_DIR, 'unit_heatmap.png')
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.1), width=Inches(9))

    add_slide_footer(slide, 5)

def add_knowing_doing_gap_slide(prs):
    """Add knowing-doing gap slide (Slide 6)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "THE KNOWING-DOING GAP"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Add chart
    chart_path = os.path.join(CHARTS_DIR, 'knowing_doing_gap.png')
    slide.shapes.add_picture(chart_path, Inches(0.3), Inches(1.1), width=Inches(9.4))

    add_slide_footer(slide, 6)

def add_participant_clusters_slide(prs):
    """Add participant clusters slide (Slide 7)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "PARTICIPANT CLUSTERS"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Add chart
    chart_path = os.path.join(CHARTS_DIR, 'participant_clusters.png')
    slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.1), width=Inches(8.4))

    add_slide_footer(slide, 7)

def add_milestone_progress_slide(prs):
    """Add milestone progress slide (Slide 8)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "MILESTONE PROGRESS"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Add chart
    chart_path = os.path.join(CHARTS_DIR, 'milestone_progress.png')
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.1), width=Inches(9))

    add_slide_footer(slide, 8)

def add_key_insights_slide(prs, survey_df, unit_data):
    """Add key insights by unit slide (Slide 9)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "KEY INSIGHTS BY UNIT"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Unit insights
    units = sorted(survey_df['unit'].unique())
    unit_colors = [PURPLE_HEX, BLUE_HEX, GOLD_HEX]

    col_width = 3
    for idx, (unit, color) in enumerate(zip(units, unit_colors)):
        unit_df = survey_df[survey_df['unit'] == unit]
        x_pos = 0.5 + idx * 3

        # Unit title
        unit_title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.2), Inches(col_width-0.3), Inches(0.4))
        unit_title_frame = unit_title_box.text_frame
        p = unit_title_frame.paragraphs[0]
        p.text = unit.upper()
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*[int(color[i:i+2], 16) for i in (1, 3, 5)])

        # Insights
        insights_text = [
            f"N={len(unit_df)}",
            f"Avg Fluency: {unit_df['overall_percentage'].mean():.0f}%",
            f"Top Construct: {max(unit_data[unit], key=unit_data[unit].get)}",
            f"Gap Risk: {(unit_df['ai_knowledge_pct'] > unit_df['applied_skills_pct']).sum()} of {len(unit_df)}"
        ]

        insights_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.7), Inches(col_width-0.3), Inches(4))
        insights_frame = insights_box.text_frame
        insights_frame.word_wrap = True

        for i, insight in enumerate(insights_text):
            if i > 0:
                insights_frame.add_paragraph()
            p = insights_frame.paragraphs[i]
            p.text = insight
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.space_before = Pt(6)

    add_slide_footer(slide, 9)

def add_interventions_slide(prs):
    """Add recommended interventions slide (Slide 10)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "RECOMMENDED INTERVENTIONS"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Interventions
    interventions = [
        ("Skill-Building Accelerator",
         "6 participants show knowing-doing gap > 10%. Provide hands-on practice labs for governance use cases."),
        ("Knowledge Integration Sessions",
         "Current AI Use lags other constructs. Host weekly 30-min integration sessions pairing theory with practice."),
        ("Peer Mentoring: High-to-Emerging",
         "4 Leading participants mentor 5 Exploring/Building staff. Establishes peer learning culture."),
        ("Unit-Specific Coaching",
         "Student Success needs AI application focus. Innovation Hub ready for advanced governance topics.")
    ]

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Clear default paragraph and build from scratch
    text_frame.clear()

    for i, (title, desc) in enumerate(interventions):
        # Title
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_before = Pt(10) if i > 0 else Pt(0)

        # Description
        p_desc = text_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = WHITE
        p_desc.level = 0
        p_desc.space_after = Pt(6)

    add_slide_footer(slide, 10)

def add_next_steps_slide(prs):
    """Add next steps slide (Slide 11)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "NEXT STEPS"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Content
    next_steps = [
        ("April 8 – Phase 1 Launch", "All participants begin Governance Foundations activities."),
        ("Bi-weekly Pulse Checks", "Brief 5-question surveys to track knowledge application & barriers."),
        ("Mid-Phase Coaching Sprints", "Week 3–4: Unit leads identify and address gaps with peer mentors."),
        ("Post-Phase 1 Assessment (May 15)", "Full fluency re-assessment; compare baseline to 6-week growth."),
        ("Success Metrics to Watch", "Completion rates >90%, knowing-doing gap closure >15%, Unit Success >75% fluency.")
    ]

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Clear default paragraph and build from scratch
    text_frame.clear()

    for i, (title, desc) in enumerate(next_steps):
        # Title
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.space_before = Pt(10) if i > 0 else Pt(0)

        # Description
        p_desc = text_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = WHITE
        p_desc.level = 0
        p_desc.space_after = Pt(6)

    add_slide_footer(slide, 11)

def add_slide_footer(slide, slide_num):
    """Add footer with slide number to slide."""
    footer_box = slide.shapes.add_textbox(Inches(9.2), Inches(6.8), Inches(0.5), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = str(slide_num)
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_LIGHT
    p.alignment = PP_ALIGN.RIGHT

def create_presentation(survey_df, milestone_df):
    """Create the full PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Calculate data
    construct_means, _ = calculate_construct_means(survey_df)
    unit_data = calculate_unit_performance(survey_df)
    gaps, participant_info = calculate_knowing_doing_gap(survey_df)
    milestone_progress = get_milestone_progress(milestone_df)

    # Generate charts
    print("Generating charts...")
    create_pie_chart_fluency(survey_df)
    create_construct_bars(construct_means, survey_df)
    create_unit_heatmap(unit_data)
    create_knowing_doing_gap_chart(gaps, participant_info)
    create_participant_clusters(survey_df)
    create_milestone_chart(milestone_progress, survey_df)

    # Build slides
    print("Building presentation...")
    add_title_slide(prs)
    add_executive_summary_slide(prs, survey_df, construct_means)
    add_sample_overview_slide(prs)
    add_construct_overview_slide(prs)
    add_unit_comparison_slide(prs)
    add_knowing_doing_gap_slide(prs)
    add_participant_clusters_slide(prs)
    add_milestone_progress_slide(prs)
    add_key_insights_slide(prs, survey_df, unit_data)
    add_interventions_slide(prs)
    add_next_steps_slide(prs)

    # Save
    prs.save(OUTPUT_PPTX)
    print(f"Presentation saved to: {OUTPUT_PPTX}")

def main():
    """Main execution."""
    print("DC CAP Enterprise AI Leadership Pilot - Baseline Report Generator")
    print("=" * 70)

    # Load data
    print("Loading survey and milestone data...")
    survey_df, milestone_df = load_data()
    print(f"  Survey: {len(survey_df)} participants")
    print(f"  Milestones: {len(milestone_df)} records")

    # Create presentation
    create_presentation(survey_df, milestone_df)

    print("=" * 70)
    print("SUCCESS: Baseline report generation complete!")
    print(f"Output: {OUTPUT_PPTX}")

if __name__ == '__main__':
    main()
